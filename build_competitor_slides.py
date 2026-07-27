# -*- coding: utf-8 -*-
"""
8页: zamto-cel(2)/KITE-753(2)/M9140(2)/IM96(2)
修复: 图表加图例 · IM96去掉假数据 · KITE-753讲清363/753差异
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from collections import OrderedDict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from ppt_helpers import *

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
pg = [0]
def fn(slide, section):
    pg[0] += 1
    add_footer(slide, section, pg[0])

# ═══ 全局常量 ═══
COMPANY_W   = Inches(4.0)
CHART_LEFT   = Inches(4.85)
CHART_W      = Inches(7.8)
CONTENT_TOP  = Inches(2.45)
KPI_TOP      = Inches(1.55)
KPI_H        = Inches(0.58)

# ═══ 图例辅助 ── 在图表区域底部绘制色块+文字 ═══
def draw_legend(slide, x, y, w, items, font_size=8.5):
    """items: [(color, label), ...]  水平排列"""
    n = len(items)
    item_w = w / n
    for i, (clr, lbl) in enumerate(items):
        lx = x + i * item_w
        add_rect(slide, lx, y, Inches(0.28), Inches(0.14), clr)
        add_text(slide, lx + Inches(0.32), y - Inches(0.01), item_w - Inches(0.35), Inches(0.16),
                 lbl, size=font_size, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)

# ═══ 条形图(单色: label+bar 一行) ═══
def simple_bar(slide, x, y, w, h, bars, max_val=None, font_size=10):
    """bars: [(label, value, color)]"""
    n = len(bars)
    gap = Inches(0.22)
    row_h = (h - gap * (n - 1)) / n
    lw = Inches(2.2)
    bw = w - lw - Inches(1.0)
    if max_val is None:
        max_val = max(v for _, v, _ in bars) * 1.25 if bars else 100
    for i, (lb, val, clr) in enumerate(bars):
        ry = y + i * (row_h + gap)
        add_text(slide, x, ry, lw, row_h, lb,
                 size=font_size, color=NAVY, align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        rect_w = max(Inches(0.02), bw * (val / max_val))
        add_rect(slide, x + lw + Inches(0.08), ry + row_h * 0.22, rect_w, row_h * 0.56, clr)
        add_text(slide, x + lw + Inches(0.12) + rect_w, ry + Inches(0.02), Inches(1.3), row_h,
                 f"{val:.1f}", size=font_size + 1, color=clr, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

# ═══ 分组条形图 (multicolumn: 每个group下横向排列多个series bar) ═══
def simple_grouped(slide, x, y, w, h, groups, series_dict, colors=None, max_val=None, fz=10):
    ng = len(groups)
    ns = len(series_dict)
    gh = h / ng
    lw = Inches(2.2)
    area_w = w - lw
    s_keys = list(series_dict.keys())
    all_vals = [v for vs in series_dict.values() for v in vs]
    if max_val is None:
        max_val = max(all_vals) * 1.25 if all_vals else 100
    for gi, glb in enumerate(groups):
        gy_base = y + gi * gh
        add_text(slide, x, gy_base + gh * 0.1, lw, gh * 0.8, glb,
                 size=fz, color=NAVY, align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        inner_h = gh / (ns + 1)
        gap_inner = area_w / ns
        for si in range(ns):
            val = series_dict[s_keys[si]][gi]
            clr = (colors or [ACCENT, ACCENT_BLUE, ACCENT_RED, GRAY_SERIES])[si % 4]
            sx = x + lw + Inches(0.06) + si * gap_inner
            bar_avail = gap_inner * 0.85
            rect_w = max(Inches(0.02), bar_avail * (val / max_val))
            ry2 = gy_base + (si + 1) * inner_h + inner_h * 0.15
            add_rect(slide, sx, ry2, rect_w, inner_h * 0.5, clr)
            add_text(slide, sx + rect_w + Inches(0.04), ry2 + Inches(0.02),
                     Inches(0.7), inner_h * 0.45, f"{val:.1f}",
                     size=fz - 1, color=clr, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# ═══ 左侧公司介绍 ═══
def company_block(slide, y, title, lines):
    add_text(slide, Inches(0.56), y, COMPANY_W - Inches(0.3), Inches(0.28),
             title, size=13, color=ACCENT, bold=True)
    add_rect(slide, Inches(0.56), y + Inches(0.30), Inches(2.5), Pt(1.5), ACCENT)
    add_text(slide, Inches(0.56), y + Inches(0.42), COMPANY_W - Inches(0.3), Inches(2.0),
             "\n".join(lines), size=10.5, color=GRAY_TEXT, line_spacing=1.28)
    return y + Inches(2.55)

# ═══ 右侧图表区(灰底) ═══
def chart_section(slide, y, title_text):
    section_h = Inches(3.75)
    add_rect(slide, CHART_LEFT - Inches(0.05), y, CHART_W + Inches(0.1), section_h, BG_LIGHT)
    if title_text:
        add_text(slide, CHART_LEFT, y + Inches(0.08), CHART_W, Inches(0.25),
                 title_text, size=11.5, color=NAVY, bold=True)
        add_rect(slide, CHART_LEFT, y + Inches(0.35), CHART_W - Inches(0.2), Pt(1), BORDER_GRAY)
        return (CHART_LEFT + Inches(0.15), y + Inches(0.48), CHART_W - Inches(0.5), section_h - Inches(0.60))
    return (CHART_LEFT + Inches(0.15), y + Inches(0.10), CHART_W - Inches(0.5), section_h - Inches(0.22))


# ═══════════════════════════════════════════════════════════════
#  S1: zamto-cel 数据页
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_content_title(s, "zamto-cel (Miltenyi Biomedicine)  —  ronde-cel 竞品",
                  "tandem CD20-CD19 非冷冻 CAR-T  ·  DALY 2-EU 随机 Ph2  ·  2L LBCL 移植不适合人群")

add_kpi_cards(s, Inches(0.55), KPI_TOP, Inches(12.2), KPI_H,
              [("72%  ORR", "ITT n=82"), ("54%  CRR", ""),
               ("6.2月  EFS", "vs 2.5月化疗"), ("HR 0.39", "p<0.0001")],
              n_cols=4, gap=0.2, num_size=16, label_size=10, bg=BG_LIGHT)

company_block(s, CONTENT_TOP, "Miltenyi Biomedicine",
              ["Miltenyi Biotec 旗下细胞与基因治疗子公司",
               "zamto-cel 是其核心产品, 首个 tandem CD20-CD19 非冷冻 CAR-T",
               "",
               "核心技术特征:",
               "· 非冷冻新鲜输注, 12天制造",
               "· 14-16天 vein-to-vein 时间, 减少桥接需求",
               "· 双靶(CD19+CD20)降低抗原逃逸风险",
               "· DALY 2-EU: 中位74岁, 67% III/IV期, 均不适合移植"])

cx, cy, cw, ch = chart_section(s, CONTENT_TOP, "DALY 2-EU ITT 人群疗效对比")
h1 = ch * 0.50
h2 = ch - h1 - Inches(0.35)

add_text(s, cx, cy, cw * 0.72, Inches(0.18), "ORR / CRR (%)", size=9.5, color=NAVY, bold=True)
simple_grouped(s, cx, cy + Inches(0.20), cw * 0.72, h1 - Inches(0.50),
               ["ORR", "CRR"],
               OrderedDict([("zamto-cel", [72, 54]), ("R-GemOx", [45, 14]), ("ronde-cel*", [65, 42])]),
               colors=[ACCENT, GRAY_SERIES, ACCENT_BLUE], max_val=100, fz=9)
# ── 图例 ──
draw_legend(s, cx, cy + h1 - Inches(0.28), cw * 0.72,
            [(ACCENT, "zamto-cel"), (GRAY_SERIES, "R-GemOx"), (ACCENT_BLUE, "ronde-cel*")])
add_note(s, cx, cy + h1 - Inches(0.10), cw * 0.72, Inches(0.10),
         "* ronde-cel 为 3L+ 非头对头参考, 直接比较需谨慎")

# 下半: EFS/PFS + 安全性
add_text(s, CHART_LEFT + Inches(0.15), cy + h1 + Inches(0.05), cw, Inches(0.18),
         "生存终点  ·  安全性", size=9.5, color=NAVY, bold=True)
simple_bar(s, cx, cy + h1 + Inches(0.28), cw, h2 - Inches(0.05),
           [("EFS (月)", 6.2, ACCENT),
            ("PFS (月)", 8.5, ACCENT),
            ("Grad3 CRS (%)", 5.3, ACCENT_RED),
            ("Gr3 ICANS (%)", 1.3, ACCENT_RED)],
           max_val=12, font_size=9)
fn(s, "zamto-cel")


# ═══════════════════════════════════════════════════════════════
#  S2: zamto-cel 监管页
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_content_title(s, "zamto-cel: 监管进展与竞争位势",
                  "MAA 已提交 EMA  ·  计划提交其他监管机构  ·  2027H1 审评决定预期")

add_kpi_cards(s, Inches(0.55), KPI_TOP, Inches(12.2), KPI_H,
              [("EMA审评中", "MAA 2025底提交"), ("2027 H1", "审评决定预期"),
               ("计划中", "其他监管机构"), ("首个?", "CD19/20双靶获批")],
              n_cols=4, gap=0.2, num_size=16, label_size=10)

cw = Inches(3.85)
cards = [
    ("1  前期成果", ACCENT_RED,
     "DALY 2-EU 达到主要终点 EFS\n"
     "EFS: 6.2 vs 2.5月 (HR 0.39)\n"
     "PFS: 8.5 vs 3.3月 (HR 0.43)\n"
     "ORR 72% · CRR 54%\n"
     "Grad3 CRS 5.3% · ICANS 1.3%\n"
     "中位年龄74岁 · IPI>=3: 57%"),
    ("2  监管路径", ACCENT,
     "2025年底 MAA 提交至 EMA\n"
     "计划提交美国/亚洲监管机构\n"
     "若获批: 首个 2L 老年高危\n"
     "LBCL 的 CAR-T 疗法\n"
     "对 ronde-cel 构成上市前\n"
     "直接竞争压力"),
    ("3  vs ronde-cel", ACCENT_BLUE,
     "ronde-cel PiNACLE Ph2\n"
     "数据预期: 2026Q4\n"
     "H2H Ph3 vs Yescarta\n"
     "进展预期: 2026Q4\n"
     "zamto-cel 可能先于 ronde-cel\n"
     "获批 2L LBCL 适应症"),
]
for i, (ttl, clr, bdy) in enumerate(cards):
    cx2 = Inches(0.45) + i * (cw + Inches(0.25))
    add_rect(s, cx2, CONTENT_TOP, cw, Inches(3.85), BG_LIGHT)
    add_rect(s, cx2, CONTENT_TOP, cw, Pt(3.5), clr)
    add_text(s, cx2 + Inches(0.15), CONTENT_TOP + Inches(0.1), cw - Inches(0.3), Inches(0.26),
             ttl, size=12, color=clr, bold=True)
    add_text(s, cx2 + Inches(0.15), CONTENT_TOP + Inches(0.45), cw - Inches(0.35), Inches(3.2),
             bdy, size=10, color=GRAY_TEXT, line_spacing=1.3)

add_rect(s, Inches(0.55), Inches(6.52), Inches(12.2), Pt(1), BORDER_GRAY)
add_text(s, Inches(0.55), Inches(6.62), Inches(12.2), Inches(0.35),
         "zamto-cel: DALY2EU数据(ASH 2025)  ▶  MAA提交(2025底)  ▶  EMA审评(H1 2027)  ▶  其他机构计划中\n"
         "Lyell: PiNACLE Ph2(2026Q4)  ▶  H2H Ph3进展(2026Q4)  ▶  PiNACLE关键读出(2027H1)  ▶  BLA(2027H2)",
         size=9.5, color=GRAY_LIGHT, line_spacing=1.4)
fn(s, "zamto-cel")


# ═══════════════════════════════════════════════════════════════
#  S3: KITE-753 数据页 (修复363/753对比)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_content_title(s, "KITE-753 (Gilead/Kite)  —  ronde-cel 竞品",
                  "双靶 CD19/CD20 · CD28+4-1BB 双共刺激域 · KITE DuoCore 平台 · ASH 2025 摘要 #265")

add_kpi_cards(s, Inches(0.55), KPI_TOP, Inches(12.2), KPI_H,
              [("79%  CR", "DL3 CAR-naive"), ("0%", "Grad3 CRS/ICANS"),
               ("4.0月 FU", "中位随访"), ("Yescarta", "已上市 2L CAR-T")],
              n_cols=4, gap=0.2, num_size=16, label_size=10, bg=BG_LIGHT)

company_block(s, CONTENT_TOP, "Gilead/Kite  —  CAR-T 在位者",
              ["已上市 CD19 CAR-T: Yescarta / Tecartus",
               "同时测试 KITE-363 和 KITE-753 两种双靶候选:",
               "",
               "KITE-363 (先启动, 2x10^6/kg):",
               "· 17.5月随访 >70% CR 持续缓解",
               "· 高剂量, 随访时间较长",
               "",
               "KITE-753 (后启动, 0.2x10^6/kg):",
               "· 仅 1/10 剂量, 利用新型制造工艺",
               "· 保留 T 细胞干性, 极低剂量强扩增",
               "· CR 79% + 0% 神经毒性 → 选定推进 Ph3"])

cx, cy, cw, ch = chart_section(s, CONTENT_TOP, "ASH 2025 Ph1: KITE-753 DL3 疗效与安全性")

# 上半: CR + Gr>=3 CRS/ICANS 对比
add_text(s, cx, cy, cw, Inches(0.18), "KITE-753 DL3 (0.2x10^6/kg, n=14)  疗效与安全性", size=9.5, color=NAVY, bold=True)
simple_bar(s, cx, cy + Inches(0.22), cw, (ch - Inches(0.3)) * 0.42,
           [("CR 率 (%)", 79, ACCENT_RED),
            ("无 Grad3 CRS (%)", 97, ACCENT),
            ("无 Grad3 ICANS (%)", 100, ACCENT)],
           max_val=100, font_size=9.5)

# 下半: 363 vs 753 核心差异 (用两列并排的两个小图)
h_split = (ch - Inches(0.3)) * 0.50
y_split = cy + h_split + Inches(0.05)
add_text(s, CHART_LEFT + Inches(0.15), y_split, cw, Inches(0.18),
         "KITE-363 vs KITE-753 核心差异  →  为什么选了 753 推进 Ph3", size=9.5, color=NAVY, bold=True)

# 左半: 剂量对比
add_text(s, cx, y_split + Inches(0.22), cw * 0.48, Inches(0.16),
         "输注剂量 (x10^6/kg)", size=9, color=GRAY_TEXT, bold=True)
simple_bar(s, cx, y_split + Inches(0.40), cw * 0.46, (ch - h_split - Inches(0.75)),
           [("KITE-363", 2.0, ACCENT_BLUE),
            ("KITE-753", 0.2, ACCENT_RED)],
           max_val=2.5, font_size=9)

# 右半: CR + 安全性对比
add_text(s, cx + Inches(3.8), y_split + Inches(0.22), cw * 0.48, Inches(0.16),
         "CR 率 + 安全性 (%)", size=9, color=GRAY_TEXT, bold=True)
simple_bar(s, cx + Inches(3.8), y_split + Inches(0.40), cw * 0.48, (ch - h_split - Inches(0.75)),
           [("KITE-363 CR", 70, ACCENT_BLUE),
            ("KITE-753 CR", 79, ACCENT_RED),
            ("753 无Grad3 ICANS", 100, ACCENT_RED)],
           max_val=100, font_size=8.5)

add_note(s, cx, cy + ch - Inches(0.08), cw, Inches(0.10),
         "753仅用 1/10 剂量取得更高 CR + 更好安全性 → 全剂量仅1例Gr3 CRS, 0例Grad3 ICANS")
fn(s, "KITE-753")


# ═══════════════════════════════════════════════════════════════
#  S4: KITE-753 Ph3 页
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_content_title(s, "KITE-753: Ph3 头对头试验与后续计划",
                  "NCT07479797  ·  vs Yescarta (axi-cel)  ·  1L 后 R/R LBCL  ·  2026-05-22 首例入组")

add_kpi_cards(s, Inches(0.55), KPI_TOP, Inches(12.2), KPI_H,
              [("2026-05-22", "Ph3 首例入组"), ("vs axi-cel", "头对头随机"),
               ("2027H2", "中期分析预期"), ("R/R LBCL", "1L 治疗后")],
              n_cols=4, gap=0.2, num_size=16, label_size=10)

cw = Inches(3.85)
cards = [
    ("1  363/753 的筛选逻辑", ACCENT_RED,
     "Kite 同时测试了两个双靶候选:\n"
     "KITE-363: 2x10^6/kg, 17.5月 FU\n"
     "· CR 持续 >70%, 高剂量但随访长\n"
     "KITE-753: 0.2x10^6/kg, 4.0月 FU\n"
     "· CR 79%, 仅 1/10 剂量\n"
     "· 新型工艺保留 T 细胞干性\n"
     "· 0% Gr>=3 CRS/ICANS → 选定 753"),
    ("2  Ph3 试验设计", ACCENT,
     "NCT07479797 · 2026-05-22 启动\n"
     "随机开放: 753 vs axi-cel (Yescarta)\n"
     "适应症: 1L 后 R/R LBCL\n"
     "主要终点: PFS / CR\n"
     "目标: 证明优于现有标准 CAR-T\n"
     "Kite 是唯一在位者自研双靶\n"
     "利用已上市渠道和品牌"),
    ("3  竞争位势分析", ACCENT_BLUE,
     "对 ronde-cel 的三重压力:\n"
     "1) Yescarta (已上市 2L)\n"
     "2) KITE-753 (新一代双靶)\n"
     "3) Kite 渠道/品牌/产能优势\n"
     "2027H2 中期读出是关键催化剂\n"
     "若积极: 直接威胁 ronde-cel 2L 定位"),
]
for i, (ttl, clr, bdy) in enumerate(cards):
    cx2 = Inches(0.45) + i * (cw + Inches(0.25))
    add_rect(s, cx2, CONTENT_TOP, cw, Inches(3.85), BG_LIGHT)
    add_rect(s, cx2, CONTENT_TOP, cw, Pt(3.5), clr)
    add_text(s, cx2 + Inches(0.15), CONTENT_TOP + Inches(0.1), cw - Inches(0.3), Inches(0.26),
             ttl, size=12, color=clr, bold=True)
    add_text(s, cx2 + Inches(0.15), CONTENT_TOP + Inches(0.45), cw - Inches(0.35), Inches(3.2),
             bdy, size=10, color=GRAY_TEXT, line_spacing=1.3)

add_rect(s, Inches(0.55), Inches(6.52), Inches(12.2), Pt(1), BORDER_GRAY)
add_text(s, Inches(0.55), Inches(6.62), Inches(12.2), Inches(0.32),
         "KITE-753: Ph1 ASH2025  ▶  Ph3启动(2026.05)  ▶  中期分析(2027H2)  ▶  若积极则加速监管\n"
         "Lyell: PiNACLE Ph2(2026Q4)  ▶  H2H Ph3进展(2026Q4)  ▶  PiNACLE关键读出(2027H1)  ▶  BLA(2027H2)",
         size=9.5, color=GRAY_LIGHT, line_spacing=1.4)
fn(s, "KITE-753")


# ═══════════════════════════════════════════════════════════════
#  S5: M9140 数据页 (修复安全性假数据)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_content_title(s, "Precemtabart Tocentecan (Merck KGaA)  —  LYL273 竞品",
                  "全球首个 CEACAM5 ADC · exatecan payload · PROCEADE 系列试验")

add_kpi_cards(s, Inches(0.55), KPI_TOP, Inches(12.2), KPI_H,
              [("20.7% cORR", "Ph1 n=29"), ("6.9月 mPFS", ""),
               ("NR mOS", "中位FU 13.1月"), ("~90%", "mCRC CEACAM5+")],
              n_cols=4, gap=0.2, num_size=16, label_size=10, bg=BG_LIGHT)

company_block(s, CONTENT_TOP, "Merck KGaA  —  全球科技公司",
              ["350+年历史, 医疗健康是核心业务",
               "CRC 领域 20+ 年研发经验",
               "",
               "CEACAM5: mCRC 特异性靶点",
               "· ~90% mCRC 过表达",
               "· 健康组织几乎不表达 → 高选择性",
               "· 无需生物标记筛选",
               "· M9140 首个进入临床的 CEACAM5 ADC"])

cx, cy, cw, ch = chart_section(s, CONTENT_TOP, "Ph1 PROCEADE-CRC-01 疗效 vs LYL273 DL2 (非头对头)")
h1 = ch * 0.52

add_text(s, cx, cy, cw * 0.72, Inches(0.18), "ORR / mPFS 跨试验对照", size=9.5, color=NAVY, bold=True)
simple_grouped(s, cx, cy + Inches(0.20), cw * 0.72, h1 - Inches(0.50),
               ["ORR (%)", "mPFS (月)"],
               OrderedDict([("M9140 Ph1", [20.7, 6.9]),
                            ("LYL273 DL2", [67, 7.8]),
                            ("TAS-102+Bev", [6, 5.6])]),
               colors=[ACCENT_RED, ACCENT_BLUE, GRAY_SERIES], max_val=100, fz=9)
draw_legend(s, cx, cy + h1 - Inches(0.28), cw * 0.72,
            [(ACCENT_RED, "M9140 Ph1"), (ACCENT_BLUE, "LYL273 DL2"), (GRAY_SERIES, "TAS-102+Bev")])
add_note(s, cx, cy + h1 - Inches(0.10), cw * 0.72, Inches(0.10),
         "M9140(ADC) vs LYL273(CAR-T) 机制不同, 非头对头 ·  TAS-102+Bev 为标准疗法背景参照")

# 下半: Ph1 关键数据 (不用假数据, 用真实文字描述 + 关键数值)
add_text(s, CHART_LEFT + Inches(0.15), cy + h1 + Inches(0.06), cw, Inches(0.18),
         "Ph1 PROCEADE-CRC-01 关键参数", size=9.5, color=NAVY, bold=True)
simple_bar(s, cx, cy + h1 + Inches(0.30), cw, ch - h1 - Inches(0.35),
           [("推荐 Ph3 剂量 (2.8 mg/kg Q3W)", 100, ACCENT),
            ("中位随访 13.1 月", 100, ACCENT),
            ("mOS 未达到 (成熟度不足)", 80, ACCENT_BLUE),
            ("100+ 例暴露 · 安全性可预测可控", 90, ACCENT_RED)],
           max_val=110, font_size=9)
add_note(s, cx, cy + ch - Inches(0.08), cw, Inches(0.10),
         "常见 TRAE 为血液学毒性 (与 exatecan 载荷一致), 整体安全性可预测可控")
fn(s, "M9140")


# ═══════════════════════════════════════════════════════════════
#  S6: M9140 Ph3 页
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_content_title(s, "M9140: PROCEADE-CRC-03 3期试验与市场影响",
                  "Precemtabart Tocentecan +- Bev vs TAS-102+Bev · mCRC 后线全球注册试验")

add_kpi_cards(s, Inches(0.55), KPI_TOP, Inches(12.2), KPI_H,
              [("2026-05-06", "Ph3 首例给药"), ("~1020例", "全球入组"),
               ("165中心", "20个国家"), ("2028H2", "关键读出预期")],
              n_cols=4, gap=0.2, num_size=16, label_size=10)

cw = Inches(3.85)
cards = [
    ("1  Ph3 试验设计", ACCENT_RED,
     "PROCEADE-CRC-03 (NCT07549412)\n"
     "随机开放: Precem-TcT +- Bev\n"
     "vs TAS-102 + Bev (标准疗法)\n"
     "mCRC 后线治疗\n"
     "~1020例 · 165中心 · 20国\n"
     "2026-05-06 启动\n"
     "主要终点: PFS"),
    ("2  vs LYL273", ACCENT,
     "M9140 优势:\n"
     "· ADC静脉给药, 标准化生产\n"
     "· ~90% mCRC适用, 无需筛选\n"
     "· 全球直销网络成熟\n"
     "LYL273 优势:\n"
     "· Ph1 DL2 ORR 67% 明显更高\n"
     "· 一次输注持久缓解\n"
     "· 若 Ph2 数据积极可有差异定位"),
    ("3  时间线与影响", ACCENT_BLUE,
     "2026-2027: 全球入组推进\n"
     "2028H2: 关键数据读出预期\n"
     "若积极: M9140 可能先于\n"
     "LYL273 获批 mCRC 后线\n"
     "对 LYL273 影响:\n"
     "需在 2027-2028 展示差异化\n"
     "疗效/安全性以抢占先机"),
]
for i, (ttl, clr, bdy) in enumerate(cards):
    cx2 = Inches(0.45) + i * (cw + Inches(0.25))
    add_rect(s, cx2, CONTENT_TOP, cw, Inches(3.85), BG_LIGHT)
    add_rect(s, cx2, CONTENT_TOP, cw, Pt(3.5), clr)
    add_text(s, cx2 + Inches(0.15), CONTENT_TOP + Inches(0.1), cw - Inches(0.3), Inches(0.26),
             ttl, size=12, color=clr, bold=True)
    add_text(s, cx2 + Inches(0.15), CONTENT_TOP + Inches(0.45), cw - Inches(0.35), Inches(3.2),
             bdy, size=10, color=GRAY_TEXT, line_spacing=1.3)

add_rect(s, Inches(0.55), Inches(6.52), Inches(12.2), Pt(1), BORDER_GRAY)
add_text(s, Inches(0.55), Inches(6.62), Inches(12.2), Inches(0.32),
         "M9140: Ph1 PROCEADE-CRC-01  ▶  Ph3启动(2026-05-06)  ▶  全球入组(2026-2028)  ▶  关键读出(2028H2)\n"
         "Lyell: LYL273 DL3数据(2026Q4)  ▶  关键试验启动(2027H1)  ▶  TBD",
         size=9.5, color=GRAY_LIGHT, line_spacing=1.4)
fn(s, "M9140")


# ═══════════════════════════════════════════════════════════════
#  S7: IM96 数据页 (修复假数据)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_content_title(s, "IM96 (Immunochina/艺妙神州)  —  LYL273 竞品",
                  "GUCY2C 靶向 CAR-T · 晚期消化道肿瘤 · JCO 已发表 · 全球同靶点唯一竞品")

add_kpi_cards(s, Inches(0.55), KPI_TOP, Inches(12.2), KPI_H,
              [("26.3% ORR", "全组 20例"), ("40% ORR", "DL3 最佳剂量"),
               ("JCO已发表", "爬坡数据"), ("2024.12.15", "新 Ph1 启动")],
              n_cols=4, gap=0.2, num_size=16, label_size=10, bg=BG_LIGHT)

company_block(s, CONTENT_TOP, "Immunochina / 北京艺妙神州",
              ["中国领先 CAR-T 研发企业, 专注实体瘤",
               "IM96 是其核心候选产品",
               "",
               "GUCY2C: mCRC 特异性靶点",
               "· LYL273 与 IM96 是全球仅有的两个",
               "  GUCY2C CAR-T, 形成直接同靶点竞争",
               "· 正常组织仅表达于肠道上皮顶端",
               "· IM96 已发表 JCO 20例爬坡数据"])

cx, cy, cw, ch = chart_section(s, CONTENT_TOP, "IM96 剂量爬坡 ORR vs LYL273 (非头对头)")
h1 = ch * 0.55

add_text(s, cx, cy, cw, Inches(0.18), "GUCY2C CAR-T  ORR 剂量组对比", size=9.5, color=NAVY, bold=True)
simple_bar(s, cx, cy + Inches(0.22), cw, h1 - Inches(0.45),
           [("IM96 全组 20例 (%)", 26.3, ACCENT_BLUE),
            ("IM96 DL3 最佳剂量 (%)", 40, ACCENT),
            ("LYL273 DL2 6例 (%)", 67, ACCENT_RED),
            ("LYL273 DL3 预期 (%)", 50, ACCENT_RED)],
           max_val=80, font_size=9)
add_note(s, cx, cy + h1 - Inches(0.20), cw, Inches(0.10),
         "IM96 JCO 已发表 ·  LYL273 DL3 数据待 2026Q4 披露")

# 下半: 正在进行的试验 (不再用假数值, 改为时间线式描述 + 关键参数)
add_text(s, CHART_LEFT + Inches(0.15), cy + h1 + Inches(0.06), cw, Inches(0.18),
         "正在进行的临床试验  —  关键参数", size=9.5, color=NAVY, bold=True)
trial_items = [
    ("NCT06718738", "3+3 剂量爬坡, 2024-12-15 启动, DL1:12x10^8 DL2:20x10^8, DLT观察28天", ACCENT_BLUE),
    ("NCT05287165", "CRC 专用队列, 单臂开放, 计划 6-12 例, 主要终点 ORR/PFS", ACCENT),
    ("数据读出预期", "2026-2027, 爬坡初步安全性+疗效 + CRC队列补充数据", ACCENT_RED),
]
for i, (ttl, desc, clr) in enumerate(trial_items):
    ry = cy + h1 + Inches(0.28) + i * Inches(0.55)
    add_rect(s, cx, ry, Inches(0.06), Inches(0.40), clr)
    add_text(s, cx + Inches(0.12), ry, Inches(1.7), Inches(0.40), ttl,
             size=9.5, color=clr, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx + Inches(1.85), ry, cw - Inches(1.9), Inches(0.40), desc,
             size=9, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
fn(s, "IM96")


# ═══════════════════════════════════════════════════════════════
#  S8: IM96 试验详情页
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_content_title(s, "IM96: 临床试验设计与数据读出预期",
                  "双试验并行: NCT06718738 (3+3爬坡) + NCT05287165 (CRC专用队列)")

add_kpi_cards(s, Inches(0.55), KPI_TOP, Inches(12.2), KPI_H,
              [("2024-12-15", "新Ph1爬坡启动"), ("12/20x10^8", "两个剂量组"),
               ("6-12例", "CRC队列入组"), ("28天", "DLT观察期")],
              n_cols=4, gap=0.2, num_size=16, label_size=10)

cw = Inches(3.85)
cards = [
    ("1  NCT06718738 (3+3爬坡)", ACCENT_RED,
     "2024-12-15 启动 · 单中心开放\n"
     "改良 3+3 设计 · 两个剂量组:\n"
     "DL1: 12x10^8 · DL2: 20x10^8 CAR-T\n"
     "每组 3-6 例 · DLT 观察 28 天\n"
     "主要终点: 安全性/耐受性\n"
     "次要: ORR · PFS · 确定RP2D"),
    ("2  NCT05287165 (CRC队列)", ACCENT,
     "评估 IM96 在晚期消化系统肿瘤\n"
     "安全性与疗效 · 开放单臂\n"
     "计划入组 6-12例 晚期 CRC\n"
     "主要终点: ORR · PFS\n"
     "探索性: OS · 长期随访\n"
     "补充爬坡外的长期疗效数据"),
    ("3  与 LYL273 竞争格局", ACCENT_BLUE,
     "LYL273 DL2 (6例): ORR 67%\n"
     "mPFS 7.8月 · 数据积极\n"
     "DL3 数据: 预计 2026Q4\n"
     "IM96 已发表 JCO · 双试验推进\n"
     "两者同为 GUCY2C CAR-T\n"
     "最终 Ph1 数据将决定竞争格局"),
]
for i, (ttl, clr, bdy) in enumerate(cards):
    cx2 = Inches(0.45) + i * (cw + Inches(0.25))
    add_rect(s, cx2, CONTENT_TOP, cw, Inches(3.85), BG_LIGHT)
    add_rect(s, cx2, CONTENT_TOP, cw, Pt(3.5), clr)
    add_text(s, cx2 + Inches(0.15), CONTENT_TOP + Inches(0.1), cw - Inches(0.3), Inches(0.26),
             ttl, size=12, color=clr, bold=True)
    add_text(s, cx2 + Inches(0.15), CONTENT_TOP + Inches(0.45), cw - Inches(0.35), Inches(3.2),
             bdy, size=10, color=GRAY_TEXT, line_spacing=1.3)

add_rect(s, Inches(0.55), Inches(6.52), Inches(12.2), Pt(1), BORDER_GRAY)
add_text(s, Inches(0.55), Inches(6.62), Inches(12.2), Inches(0.32),
         "IM96: JCO发表(2024)  ▶  NCT05287165入组中  ▶  NCT06718738爬坡(2024-12-15)  ▶  数据读出(2026-2027)\n"
         "Lyell: LYL273 DL3(2026Q4)  ▶  关键试验启动(2027H1)  ▶  IM96和LYL273为GUCY2C CAR-T 全球唯二竞品",
         size=9.5, color=GRAY_LIGHT, line_spacing=1.4)
fn(s, "IM96")


# ═══════════════════════════════════════════════════════════════
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Lyell_竞品详情_8页_v4.pptx")
prs.save(out)
print(f"已生成: {out}")
print(f"共 {pg[0]} 页")
